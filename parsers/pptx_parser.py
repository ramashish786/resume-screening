"""
parsers/pptx_parser.py
───────────────────────
Extracts text from .pptx resume/portfolio files using python-pptx.
Iterates over all slides, text frames, and table cells.
"""

from __future__ import annotations

import hashlib
import io

from loguru import logger

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    Presentation = None  # type: ignore


def parse_pptx(file_bytes: bytes, file_name: str) -> dict:
    """
    Extract text from a .pptx file.

    Args:
        file_bytes: Raw bytes of the uploaded PPTX.
        file_name:  Original filename.

    Returns:
        dict with keys: raw_text, page_count, word_count, warnings, file_hash
    """
    if Presentation is None:
        raise ImportError("python-pptx is not installed. Run: pip install python-pptx")

    warnings: list[str] = []
    file_hash = hashlib.md5(file_bytes).hexdigest()

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        raise ValueError(f"Could not open PPTX '{file_name}': {e}") from e

    slide_texts: list[str] = []
    page_count = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            # Text frames (most common in resumes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        slide_parts.append(line)

            # Tables (skills grids, experience tables)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]
                    if cells:
                        slide_parts.append(" | ".join(cells))

        if slide_parts:
            slide_texts.append(f"[Slide {slide_idx}]\n" + "\n".join(slide_parts))
        else:
            warnings.append(f"Slide {slide_idx}: no text found (may be image-only)")

    raw_text = "\n\n".join(slide_texts)
    word_count = len(raw_text.split())

    if word_count < 30:
        warnings.append(
            "Very little text extracted from PPTX. "
            "Most slides may be image-based — text may not be fully captured."
        )

    logger.info(
        f"PPTX parsed: {file_name} | slides={page_count} | words={word_count} | warnings={len(warnings)}"
    )

    return {
        "raw_text": raw_text,
        "page_count": page_count,
        "word_count": word_count,
        "warnings": warnings,
        "file_hash": file_hash,
    }